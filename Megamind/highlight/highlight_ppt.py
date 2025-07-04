from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw

def export_slide_as_image(pptx_path, slide_number, output_image):
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_number-1]
    # Export as PNG (requires comtypes on Windows, otherwise, use screenshot workaround)
    # For cross-platform, render blank slide and draw text if needed
    width = prs.slide_width
    height = prs.slide_height
    img = Image.new('RGB', (int(width/9525), int(height/9525)), color='white')
    draw = ImageDraw.Draw(img)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            draw.text((50, 50), paragraph.text, fill='black')
    img.save(output_image)
    return output_image
