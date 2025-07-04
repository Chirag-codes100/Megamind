from pptx import Presentation

def extract_ppt_text_by_slide(ppt_path):
    prs = Presentation(ppt_path)
    slides_text = {}
    for i, slide in enumerate(prs.slides, start=1):
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
        slides_text[i] = "\n".join(text)
    return slides_text
